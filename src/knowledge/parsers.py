"""Document Parsers for Phase 3 - all file types."""

import csv
import html
import io
import json
import logging
import os
import re
import xml.etree.ElementTree as ET
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional
from datetime import datetime

logger = logging.getLogger(__name__)


@dataclass
class ParsedDocument:
    text: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    pages: List[str] = field(default_factory=list)
    tables: List[List[List[str]]] = field(default_factory=list)
    images: List[str] = field(default_factory=list)
    headings: List[str] = field(default_factory=list)
    language: str = "en"
    word_count: int = 0
    char_count: int = 0


class DocumentParser(ABC):
    @abstractmethod
    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument: ...

    @abstractmethod
    def supports(self, extension: str) -> bool: ...


class TXTParser(DocumentParser):
    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument:
        text = content.decode('utf-8', errors='replace')
        return ParsedDocument(
            text=text,
            word_count=len(text.split()),
            char_count=len(text),
            metadata={"format": "text", "encoding": "utf-8"},
        )

    def supports(self, extension: str) -> bool:
        return extension in ('.txt', '.text', '.log', '.cfg', '.ini', '.conf')


class MDParser(DocumentParser):
    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument:
        text = content.decode('utf-8', errors='replace')
        headings = re.findall(r'^(#{1,6})\s+(.+)$', text, re.MULTILINE)
        clean_text = re.sub(r'[#*`~\[\]()>|_-]', ' ', text)
        clean_text = re.sub(r'\s+', ' ', clean_text).strip()
        return ParsedDocument(
            text=clean_text,
            headings=[h[1] for h in headings],
            word_count=len(clean_text.split()),
            char_count=len(clean_text),
            metadata={"format": "markdown", "headings": len(headings)},
        )

    def supports(self, extension: str) -> bool:
        return extension in ('.md', '.markdown', '.mdown', '.mdx')


class JSONParser(DocumentParser):
    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument:
        try:
            data = json.loads(content.decode('utf-8', errors='replace'))
            text = json.dumps(data, indent=2, default=str)
        except Exception:
            text = content.decode('utf-8', errors='replace')
        return ParsedDocument(
            text=text,
            word_count=len(text.split()),
            char_count=len(text),
            metadata={"format": "json"},
        )

    def supports(self, extension: str) -> bool:
        return extension == '.json'


class XMLParser(DocumentParser):
    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument:
        try:
            root = ET.fromstring(content)
            text = ET.tostring(root, encoding='unicode', method='text')
            text = re.sub(r'\s+', ' ', text).strip()
        except Exception:
            text = content.decode('utf-8', errors='replace')
        return ParsedDocument(
            text=text,
            word_count=len(text.split()),
            char_count=len(text),
            metadata={"format": "xml"},
        )

    def supports(self, extension: str) -> bool:
        return extension in ('.xml', '.xsd', '.xsl', '.xslt', '.svg')


class HTMLParser(DocumentParser):
    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument:
        text = content.decode('utf-8', errors='replace')
        clean = re.sub(r'<[^>]+>', ' ', text)
        clean = html.unescape(clean)
        clean = re.sub(r'\s+', ' ', clean).strip()
        title_match = re.search(r'<title[^>]*>(.*?)</title>', text, re.IGNORECASE | re.DOTALL)
        title = title_match.group(1).strip() if title_match else ""
        headings = re.findall(r'<h[1-6][^>]*>(.*?)</h[1-6]>', text, re.IGNORECASE | re.DOTALL)
        return ParsedDocument(
            text=clean,
            headings=[h.strip() for h in headings],
            word_count=len(clean.split()),
            char_count=len(clean),
            metadata={"format": "html", "title": title},
        )

    def supports(self, extension: str) -> bool:
        return extension in ('.html', '.htm', '.xhtml')


class CSVParser(DocumentParser):
    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument:
        try:
            text_content = content.decode('utf-8', errors='replace')
            reader = csv.reader(io.StringIO(text_content))
            rows = []
            for row in reader:
                rows.append(row)
            text = '\n'.join([','.join(row) for row in rows])
            return ParsedDocument(
                text=text,
                tables=[rows],
                word_count=len(text.split()),
                char_count=len(text),
                metadata={"format": "csv", "rows": len(rows), "columns": len(rows[0]) if rows else 0},
            )
        except Exception as e:
            logger.error(f"CSV parse error: {e}")
            text = content.decode('utf-8', errors='replace')
            return ParsedDocument(text=text, metadata={"format": "csv"})

    def supports(self, extension: str) -> bool:
        return extension == '.csv'


class PDFParser(DocumentParser):
    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument:
        pages = []
        text = ""
        try:
            import fitz
            doc = fitz.open(stream=content, filetype="pdf")
            for i, page in enumerate(doc):
                page_text = page.get_text()
                pages.append(page_text)
                text += page_text + "\n\n"
            doc.close()
        except ImportError:
            logger.warning("PyMuPDF not available, trying pypdf")
            try:
                from PyPDF2 import PdfReader
                import io
                reader = PdfReader(io.BytesIO(content))
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    pages.append(page_text)
                    text += page_text + "\n\n"
            except ImportError:
                logger.error("No PDF library available (need PyMuPDF or PyPDF2)")
                text = content.decode('utf-8', errors='replace')
                pages = [text]
        except Exception as e:
            logger.error(f"PDF parse error: {e}")
            text = content.decode('utf-8', errors='replace')
            pages = [text]

        clean_text = re.sub(r'\s+', ' ', text).strip()
        return ParsedDocument(
            text=clean_text,
            pages=pages,
            word_count=len(clean_text.split()),
            char_count=len(clean_text),
            metadata={"format": "pdf", "pages": len(pages)},
        )

    def supports(self, extension: str) -> bool:
        return extension == '.pdf'


class DOCXParser(DocumentParser):
    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument:
        text = ""
        try:
            from docx import Document
            import io
            doc = Document(io.BytesIO(content))
            paragraphs = []
            for para in doc.paragraphs:
                paragraphs.append(para.text)
            text = '\n'.join(paragraphs)
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    table_data.append([cell.text for cell in row.cells])
                tables.append(table_data)
        except ImportError:
            logger.error("python-docx not available")
            text = content.decode('utf-8', errors='replace')
        except Exception as e:
            logger.error(f"DOCX parse error: {e}")
            text = content.decode('utf-8', errors='replace')

        clean_text = re.sub(r'\s+', ' ', text).strip()
        return ParsedDocument(
            text=clean_text,
            tables=tables if tables else [],
            word_count=len(clean_text.split()),
            char_count=len(clean_text),
            metadata={"format": "docx"},
        )

    def supports(self, extension: str) -> bool:
        return extension in ('.docx', '.doc')


class XLSXParser(DocumentParser):
    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument:
        text = ""
        tables = []
        try:
            import openpyxl
            import io
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_data = []
                sheet_text = []
                for row in ws.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    sheet_data.append(row_data)
                    sheet_text.append('\t'.join(row_data))
                tables.append(sheet_data)
                text += f"\n--- Sheet: {sheet_name} ---\n" + '\n'.join(sheet_text)
            wb.close()
        except ImportError:
            logger.error("openpyxl not available")
            text = content.decode('utf-8', errors='replace')
        except Exception as e:
            logger.error(f"XLSX parse error: {e}")
            text = content.decode('utf-8', errors='replace')

        clean_text = re.sub(r'\s+', ' ', text).strip()
        return ParsedDocument(
            text=clean_text,
            tables=tables,
            word_count=len(clean_text.split()),
            char_count=len(clean_text),
            metadata={"format": "xlsx", "sheets": len(tables)},
        )

    def supports(self, extension: str) -> bool:
        return extension in ('.xlsx', '.xls')


class PPTXParser(DocumentParser):
    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument:
        text = ""
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(content))
            slides = []
            for slide in prs.slides:
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                slides.append('\n'.join(slide_texts))
            text = '\n\n'.join(slides)
        except ImportError:
            logger.error("python-pptx not available")
            text = content.decode('utf-8', errors='replace')
        except Exception as e:
            logger.error(f"PPTX parse error: {e}")
            text = content.decode('utf-8', errors='replace')

        clean_text = re.sub(r'\s+', ' ', text).strip()
        return ParsedDocument(
            text=clean_text,
            word_count=len(clean_text.split()),
            char_count=len(clean_text),
            metadata={"format": "pptx", "slides": len(slides) if 'slides' in dir() else 0},
        )

    def supports(self, extension: str) -> bool:
        return extension in ('.pptx', '.ppt')


class CodeParser(DocumentParser):
    LANGUAGES = {
        '.py': 'python', '.js': 'javascript', '.ts': 'typescript', '.tsx': 'typescript',
        '.java': 'java', '.cpp': 'cpp', '.c': 'c', '.h': 'c', '.hpp': 'cpp',
        '.cs': 'csharp', '.go': 'go', '.rs': 'rust', '.rb': 'ruby', '.php': 'php',
        '.swift': 'swift', '.kt': 'kotlin', '.scala': 'scala', '.sql': 'sql',
        '.sh': 'bash', '.bash': 'bash', '.zsh': 'bash', '.ps1': 'powershell',
        '.r': 'r', '.m': 'matlab', '.pl': 'perl', '.lua': 'lua',
        '.dart': 'dart', '.groovy': 'groovy', '.vue': 'vue', '.svelte': 'svelte',
    }

    def parse(self, content: bytes, file_path: str = "") -> ParsedDocument:
        text = content.decode('utf-8', errors='replace')
        ext = os.path.splitext(file_path)[1].lower()
        lang = self.LANGUAGES.get(ext, 'text')

        functions = re.findall(r'(?:def|function|async\s+function|public|private|static|export\s+(?:default\s+)?(?:function|class))\s+(\w+)\s*\(', text)
        classes = re.findall(r'(?:class|struct|interface|trait|enum)\s+(\w+)', text)
        imports = re.findall(r'(?:import|from|require|include|using)\s+([\'\"\w\._/-]+)', text)

        clean_text = re.sub(r'^\s*#.*$', '', text, flags=re.MULTILINE)
        clean_text = re.sub(r'^\s*//.*$', '', clean_text, flags=re.MULTILINE)
        clean_text = re.sub(r'/\*.*?\*/', '', clean_text, flags=re.DOTALL)
        clean_text = re.sub(r'\s+', ' ', clean_text).strip()

        return ParsedDocument(
            text=text,
            word_count=len(clean_text.split()),
            char_count=len(text),
            metadata={
                "format": "code",
                "language": lang,
                "functions": functions,
                "classes": classes,
                "imports": imports,
                "lines": text.count('\n') + 1,
            },
        )

    def supports(self, extension: str) -> bool:
        return extension in self.LANGUAGES


_parsers: List[DocumentParser] = [
    PDFParser(), DOCXParser(), TXTParser(), MDParser(),
    CSVParser(), XLSXParser(), PPTXParser(),
    JSONParser(), XMLParser(), HTMLParser(),
    CodeParser(),
]


def get_parser_for_file(file_path: str) -> Optional[DocumentParser]:
    ext = os.path.splitext(file_path)[1].lower()
    for parser in _parsers:
        if parser.supports(ext):
            return parser
    return TXTParser()


def get_all_parsers() -> List[DocumentParser]:
    return _parsers
